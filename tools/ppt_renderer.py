
from __future__ import annotations

import logging
import re
import tempfile
from datetime import datetime
from pathlib import Path

from jinja2 import Environment, FileSystemLoader

from config import ReportGenerateError
from core.data_context import DataContext
from models.report_content import ReportContent, ReportSection, ReportTable

logger = logging.getLogger(__name__)

# 模板目录
_TEMPLATE_DIR = Path(__file__).parent.parent / "templates" / "ppt"
_THEMES_DIR = _TEMPLATE_DIR / "themes"

# PPT 尺寸 (16:9)
_SLIDE_WIDTH_CM = 33.867
_SLIDE_HEIGHT_CM = 19.05
_VIEWPORT_WIDTH = 1280
_VIEWPORT_HEIGHT = 720


# ---------------------------------------------------------------------------
# 分页逻辑
# ---------------------------------------------------------------------------


class _SlidePage:
    """一页 PPT 的数据。"""

    def __init__(self, template_name: str, **kwargs: object) -> None:
        self.template_name = template_name
        self.data = kwargs


def _paginate(report: ReportContent) -> list[_SlidePage]:
    """将 ReportContent 分页：封面 + 内容页 + 结论页。"""
    pages: list[_SlidePage] = []

    # 封面页
    pages.append(_SlidePage("cover.html.j2", title=report.title, date=datetime.now().strftime("%Y年%m月%d日")))

    # 内容页：每个 level=1 的 section 独占一页，level=2/3 合并到上一个 level=1
    current_sections: list[ReportSection] = []
    current_heading = ""

    for section in report.sections:
        if section.level == 1:
            # 先保存之前的页
            if current_sections:
                pages.append(_SlidePage("slide.html.j2", heading=current_heading, sections=current_sections, title=report.title))
            current_heading = section.heading
            current_sections = [section]
        else:
            # level 2/3 追加到当前页
            current_sections.append(section)

    # 最后一页
    if current_sections:
        pages.append(_SlidePage("slide.html.j2", heading=current_heading, sections=current_sections, title=report.title))

    # 结论页
    if report.conclusion:
        pages.append(_SlidePage("conclusion.html.j2", title=report.title, conclusion=report.conclusion))

    return pages


# ---------------------------------------------------------------------------
# Jinja2 渲染
# ---------------------------------------------------------------------------


def _render_html_pages(pages: list[_SlidePage], theme_name: str) -> list[str]:
    """用 Jinja2 渲染每页的 HTML 字符串。"""
    # 确定主题 CSS 路径
    theme_css = _THEMES_DIR / f"{theme_name}.css"
    if not theme_css.exists():
        logger.warning("主题 '%s' 不存在，使用默认 blue 主题", theme_name)
        theme_css = _THEMES_DIR / "blue.css"

    env = Environment(loader=FileSystemLoader(str(_TEMPLATE_DIR)), autoescape=True)
    html_pages: list[str] = []

    for page in pages:
        template = env.get_template(page.template_name)
        html = template.render(**page.data, theme_css_path=str(theme_css))
        html_pages.append(html)

    return html_pages


# ---------------------------------------------------------------------------
# Playwright 截图
# ---------------------------------------------------------------------------


async def _screenshot_pages(html_pages: list[str]) -> list[Path]:
    """用 Playwright 逐页截图，返回 PNG 文件路径列表。"""
    try:
        from playwright.async_api import async_playwright
    except ImportError as e:
        raise ReportGenerateError("ppt", detail="playwright not installed, run: uv add playwright && playwright install chromium") from e

    tmp_dir = Path(tempfile.mkdtemp(prefix="bi_ppt_"))
    images: list[Path] = []

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=True)
        page = await browser.new_page(viewport={"width": _VIEWPORT_WIDTH, "height": _VIEWPORT_HEIGHT})

        for i, html in enumerate(html_pages):
            # 写入临时 HTML 文件
            html_path = tmp_dir / f"slide_{i}.html"
            html_path.write_text(html, encoding="utf-8")

            await page.goto(f"file:///{str(html_path).replace(chr(92), '/')}")
            img_path = tmp_dir / f"slide_{i}.png"
            await page.screenshot(path=str(img_path), full_page=False)
            images.append(img_path)

        await browser.close()

    return images


# ---------------------------------------------------------------------------
# python-pptx 封装
# ---------------------------------------------------------------------------


def _build_pptx(images: list[Path], title: str, output_dir: Path) -> str:
    """将截图列表封装为 .pptx 文件。"""
    try:
        from pptx import Presentation
        from pptx.util import Cm
    except ImportError as e:
        raise ReportGenerateError("ppt", detail="python-pptx not installed, run: uv add python-pptx") from e

    prs = Presentation()
    prs.slide_width = Cm(_SLIDE_WIDTH_CM)
    prs.slide_height = Cm(_SLIDE_HEIGHT_CM)

    blank_layout = prs.slide_layouts[6]  # 空白布局

    for img_path in images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    from tools.report_generate import _safe_filename

    filename = f"{_safe_filename(title)}.pptx"
    output_path = output_dir / filename
    prs.save(str(output_path))
    return f"PPT报告已生成: {output_path}"


# ---------------------------------------------------------------------------
# PDF 导出
# ---------------------------------------------------------------------------


async def _export_pdf(html_pages: list[str], title: str, output_dir: Path) -> str:
    """用 Playwright 将所有页面合并导出为 PDF。"""
    try:
        from playwright.async_api import async_playwright
    except ImportError as e:
        raise ReportGenerateError("pdf", detail="playwright not installed, run: uv add playwright && playwright install chromium") from e

    # 合并所有页面为一个 HTML，用分页符分隔
    combined_parts: list[str] = []
    for html in html_pages:
        # 提取 <body> 内容
        body_start = html.find("<body>")
        body_end = html.find("</body>")
        if body_start >= 0 and body_end >= 0:
            content = html[body_start + 6 : body_end]
        else:
            content = html
        combined_parts.append(f'<div style="page-break-after: always;">{content}</div>')

    # 提取第一页的 <head> 作为公共 head
    head_start = html_pages[0].find("<head>")
    head_end = html_pages[0].find("</head>")
    head_content = html_pages[0][head_start + 6 : head_end] if head_start >= 0 and head_end >= 0 else ""

    combined_html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>{head_content}
<style>@page {{ size: 1280px 720px; margin: 0; }}</style>
</head>
<body>
{"".join(combined_parts)}
</body>
</html>"""

    tmp_dir = Path(tempfile.mkdtemp(prefix="bi_pdf_"))
    html_path = tmp_dir / "report.html"
    html_path.write_text(combined_html, encoding="utf-8")

    from tools.report_generate import _safe_filename

    filename = f"{_safe_filename(title)}.pdf"
    output_path = output_dir / filename

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=True)
        page = await browser.new_page(viewport={"width": _VIEWPORT_WIDTH, "height": _VIEWPORT_HEIGHT})
        await page.goto(f"file:///{str(html_path).replace(chr(92), '/')}")
        await page.pdf(
            path=str(output_path),
            width=f"{_VIEWPORT_WIDTH}px",
            height=f"{_VIEWPORT_HEIGHT}px",
            print_background=True,
        )
        await browser.close()

    return f"PDF报告已生成: {output_path}"


# ---------------------------------------------------------------------------
# CSS 主题解析 → python-pptx 原生元素（可编辑 PPT）
# ---------------------------------------------------------------------------


def _parse_theme(theme_name: str) -> dict[str, str]:
    """从 CSS 主题文件解析颜色变量，返回 dict。"""
    theme_css = _THEMES_DIR / f"{theme_name}.css"
    if not theme_css.exists():
        theme_css = _THEMES_DIR / "blue.css"

    content = theme_css.read_text(encoding="utf-8")
    theme: dict[str, str] = {}
    for m in re.finditer(r"--([\w-]+)\s*:\s*([^;]+);", content):
        theme[m.group(1)] = m.group(2).strip()

    # 解析 body background
    bg_match = re.search(r"body\s*\{[^}]*background:\s*([^;]+);", content)
    if bg_match:
        theme["body_background"] = bg_match.group(1).strip()

    # 解析封面 overlay 渐变主色
    overlay_match = re.search(r"cover-overlay\s*\{[^}]*background:\s*linear-gradient\([^,]+,\s*rgba\(([^)]+)\)", content)
    if overlay_match:
        theme["cover_gradient_start"] = overlay_match.group(1).strip()

    return theme


def _hex_to_rgbcolor(hex_color: str):
    """将 #RRGGBB 转为 pptx.dml.color.RGBColor。"""
    from pptx.dml.color import RGBColor

    hex_color = hex_color.strip().lstrip("#")
    if len(hex_color) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _is_dark_theme(theme: dict[str, str]) -> bool:
    """判断是否为深色主题。"""
    bg = theme.get("body_background", "#ffffff")
    # 解析 hex 颜色，计算亮度
    bg = bg.strip().lstrip("#")
    if len(bg) != 6:
        return False
    r, g, b = int(bg[0:2], 16), int(bg[2:4], 16), int(bg[4:6], 16)
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return luminance < 128


def _build_native_pptx(report: ReportContent, output_dir: Path, data_context: DataContext | None = None) -> str:
    """用 python-pptx 原生元素构建可编辑 PPT，视觉风格来自 CSS 主题。"""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Cm, Emu, Pt
    except ImportError as e:
        raise ReportGenerateError("ppt", detail="python-pptx not installed, run: uv add python-pptx") from e

    theme = _parse_theme(report.template_name)
    is_dark = _is_dark_theme(theme)

    accent = _hex_to_rgbcolor(theme.get("accent-color", "#0078d4"))
    title_color = _hex_to_rgbcolor(theme.get("title-color", "#1a1a2e"))
    heading_color = _hex_to_rgbcolor(theme.get("heading-color", "#0078d4"))
    text_color = _hex_to_rgbcolor(theme.get("text-color", "#333333"))
    table_header_bg = _hex_to_rgbcolor(theme.get("table-header-bg", "#0078d4"))
    table_header_color = _hex_to_rgbcolor(theme.get("table-header-color", "#ffffff"))
    table_stripe_bg = _hex_to_rgbcolor(theme.get("table-stripe-bg", "#f5f7fa"))
    caption_color = _hex_to_rgbcolor(theme.get("caption-color", "#888888"))
    footer_color = _hex_to_rgbcolor(theme.get("footer-color", "#999999"))
    conclusion_bg = _hex_to_rgbcolor(theme.get("conclusion-bg", "#f0f7ff"))

    # 深色主题修正：确保文字在白色背景上可读
    if is_dark:
        # 原生 PPT 默认白色背景，深色主题的文字颜色需要反转
        text_color = RGBColor(0x1A, 0x1A, 0x2E)       # 深色正文
        title_color = RGBColor(0x0A, 0x0A, 0x1E)       # 深色标题
        heading_color = accent                          # 用强调色做标题
        caption_color = RGBColor(0x66, 0x66, 0x66)
        footer_color = RGBColor(0x99, 0x99, 0x99)
        # 表格也修正
        table_header_color = RGBColor(0xFF, 0xFF, 0xFF)
        table_stripe_bg = RGBColor(0xF0, 0xF4, 0xF8)
        conclusion_bg = RGBColor(0xE8, 0xF4, 0xFD)

    prs = Presentation()
    prs.slide_width = Cm(_SLIDE_WIDTH_CM)
    prs.slide_height = Cm(_SLIDE_HEIGHT_CM)
    blank_layout = prs.slide_layouts[6]

    # 布局常量 (cm)
    margin_left = Cm(1.8)
    margin_right = Cm(1.8)
    margin_top = Cm(1.2)
    content_width = prs.slide_width - margin_left - margin_right
    # 内容区底部边界（留出页脚空间）
    content_bottom = Cm(17.8)
    footer_y = Cm(18.2)

    # ── 封面页 ──
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = accent

    # 装饰条 — 底部细线
    deco_line = slide.shapes.add_shape(1, Cm(0), Cm(17.5), prs.slide_width, Cm(0.06))
    deco_line.fill.solid()
    deco_line.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    deco_line.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Cm(3), Cm(5.5), prs.slide_width - Cm(6), Cm(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = report.title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)

    # 日期 + 作者
    p2 = tf.add_paragraph()
    p2.text = datetime.now().strftime("%Y年%m月%d日") + "  智问BI 自动生成"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    p2.alignment = PP_ALIGN.CENTER

    # ── 内容页 ──
    pages = _paginate(report)
    for page in pages:
        if page.template_name != "slide.html.j2":
            continue

        slide = prs.slides.add_slide(blank_layout)

        # 左侧装饰色条 — 区分主次的视觉锚点
        accent_bar = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(0.4), prs.slide_height)
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()

        # 页面标题
        title_bar = slide.shapes.add_textbox(margin_left, margin_top, content_width, Cm(1.4))
        tf = title_bar.text_frame
        p = tf.paragraphs[0]
        p.text = str(page.data.get("heading", ""))
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = title_color

        # 标题下装饰线 — 渐变效果用两段模拟
        line_full = slide.shapes.add_shape(1, margin_left, margin_top + Cm(1.5), content_width, Cm(0.08))
        line_full.fill.solid()
        line_full.fill.fore_color.rgb = accent
        line_full.line.fill.background()

        # 内容区域 — 从标题下方开始
        y_cursor = margin_top + Cm(2.0)
        indent = Cm(0.4)

        sections = page.data.get("sections", [])
        for section in sections:
            if not isinstance(section, ReportSection):
                continue
            if y_cursor >= content_bottom:
                break

            # 二/三级标题
            if section.level > 1 and section.heading:
                font_size = Pt(20) if section.level == 2 else Pt(17)
                # 估算高度
                est_h = Cm(1.1)
                h_box = slide.shapes.add_textbox(
                    margin_left + indent, y_cursor, content_width - indent, est_h
                )
                h_tf = h_box.text_frame
                h_tf.word_wrap = True
                h_p = h_tf.paragraphs[0]
                h_p.text = section.heading
                h_p.font.size = font_size
                h_p.font.bold = True
                h_p.font.color.rgb = heading_color
                h_p.space_after = Pt(6)
                y_cursor += est_h + Cm(0.2)

            # 段落
            for para_text in section.paragraphs:
                if y_cursor >= content_bottom:
                    break
                # 估算行数：16pt 在 30cm 宽度内约 28 字/行
                char_per_line = 28
                num_lines = max(1, (len(para_text) + char_per_line - 1) // char_per_line)
                est_h = Cm(0.75) * num_lines + Cm(0.2)
                # 不能超出底部
                max_h = content_bottom - y_cursor
                if est_h > max_h:
                    est_h = max_h

                p_box = slide.shapes.add_textbox(
                    margin_left + indent, y_cursor, content_width - indent, est_h
                )
                p_tf = p_box.text_frame
                p_tf.word_wrap = True
                p_p = p_tf.paragraphs[0]
                p_p.text = para_text
                p_p.font.size = Pt(16)
                p_p.font.color.rgb = text_color
                p_p.line_spacing = Pt(26)
                p_p.space_after = Pt(6)
                y_cursor += est_h + Cm(0.15)

            # 表格
            for table in section.tables:
                if not isinstance(table, ReportTable):
                    continue
                num_cols = len(table.headers)
                if num_cols == 0:
                    continue
                # 计算可显示行数
                max_data_rows = len(table.rows)
                row_h = Cm(0.85)
                available = content_bottom - y_cursor - Cm(0.6)  # 留 caption 空间
                max_rows_fit = max(1, int(available / row_h) - 1)  # -1 表头
                display_rows = min(max_data_rows, max_rows_fit)
                total_rows = display_rows + 1  # +1 表头
                table_height = row_h * total_rows

                if y_cursor + table_height > content_bottom:
                    break

                tbl_shape = slide.shapes.add_table(
                    total_rows,
                    num_cols,
                    margin_left + indent,
                    y_cursor,
                    content_width - indent,
                    table_height,
                )
                tbl = tbl_shape.table

                # 表头
                for i, header in enumerate(table.headers):
                    cell = tbl.cell(0, i)
                    cell.text = header
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(13)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = table_header_color
                        paragraph.alignment = PP_ALIGN.CENTER
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = table_header_bg
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 数据行
                for r_idx in range(display_rows):
                    row_data = table.rows[r_idx]
                    for c_idx, cell_text in enumerate(row_data[:num_cols]):
                        cell = tbl.cell(r_idx + 1, c_idx)
                        cell.text = cell_text
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(13)
                            paragraph.font.color.rgb = text_color
                        if r_idx % 2 == 1:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = table_stripe_bg
                        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                y_cursor += table_height + Cm(0.2)

                # 表格说明
                if table.caption:
                    cap_box = slide.shapes.add_textbox(
                        margin_left + indent, y_cursor, content_width - indent, Cm(0.5)
                    )
                    cap_tf = cap_box.text_frame
                    cap_p = cap_tf.paragraphs[0]
                    cap_p.text = table.caption
                    cap_p.font.size = Pt(11)
                    cap_p.font.italic = True
                    cap_p.font.color.rgb = caption_color
                    y_cursor += Cm(0.55)

            # 图表
            for chart_key in section.chart_keys:
                if y_cursor >= content_bottom:
                    break
                if data_context is None:
                    break
                artifact = data_context.get_chart(chart_key)
                if artifact is None or not artifact.png_path:
                    continue
                img_path = Path(artifact.png_path)
                if not img_path.exists():
                    continue
                chart_height = Cm(7)
                chart_width = content_width * 0.9
                if y_cursor + chart_height > content_bottom:
                    break
                slide.shapes.add_picture(
                    str(img_path),
                    margin_left + indent,
                    y_cursor,
                    width=chart_width,
                    height=chart_height,
                )
                y_cursor += chart_height + Cm(0.3)

        # 页脚
        footer_box = slide.shapes.add_textbox(margin_left, footer_y, content_width, Cm(0.5))
        footer_tf = footer_box.text_frame
        footer_p = footer_tf.paragraphs[0]
        footer_p.text = report.title
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = footer_color
        footer_p.alignment = PP_ALIGN.RIGHT

    # ── 结论页 ──
    if report.conclusion:
        slide = prs.slides.add_slide(blank_layout)

        # 左侧装饰色条
        accent_bar = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(0.4), prs.slide_height)
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()

        # 标题
        title_bar = slide.shapes.add_textbox(margin_left, margin_top, content_width, Cm(1.4))
        tf = title_bar.text_frame
        p = tf.paragraphs[0]
        p.text = "结论"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = title_color

        # 标题下装饰线
        line = slide.shapes.add_shape(1, margin_left, margin_top + Cm(1.5), content_width, Cm(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = accent
        line.line.fill.background()

        # 结论背景框 — 先添加（底层）
        conc_top = Cm(4.0)
        # 估算结论文字高度
        char_per_line = 25
        num_lines = max(1, (len(report.conclusion) + char_per_line - 1) // char_per_line)
        conc_text_h = Cm(0.9) * num_lines + Cm(1.0)  # 文字高度 + 内边距
        conc_text_h = min(conc_text_h, Cm(10))  # 上限
        bg_h = conc_text_h + Cm(1.2)  # 背景框比文字多留边距

        bg_shape = slide.shapes.add_shape(1, margin_left, conc_top, content_width, bg_h)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = conclusion_bg
        bg_shape.line.fill.background()
        # 圆角效果 — 通过设置边框颜色与填充一致模拟
        # 移到最底层
        sp = bg_shape._element
        sp.getparent().remove(sp)
        slide.shapes._spTree.insert(2, sp)

        # 左侧竖线装饰 — 结论框内
        conc_accent = slide.shapes.add_shape(1, margin_left + Cm(0.3), conc_top + Cm(0.3), Cm(0.12), bg_h - Cm(0.6))
        conc_accent.fill.solid()
        conc_accent.fill.fore_color.rgb = accent
        conc_accent.line.fill.background()

        # 结论内容
        conc_box = slide.shapes.add_textbox(
            margin_left + Cm(1.0), conc_top + Cm(0.6), content_width - Cm(1.6), conc_text_h
        )
        conc_tf = conc_box.text_frame
        conc_tf.word_wrap = True
        conc_p = conc_tf.paragraphs[0]
        conc_p.text = report.conclusion
        conc_p.font.size = Pt(18)
        conc_p.font.color.rgb = text_color
        conc_p.line_spacing = Pt(32)

        # 页脚
        footer_box = slide.shapes.add_textbox(margin_left, footer_y, content_width, Cm(0.5))
        footer_tf = footer_box.text_frame
        footer_p = footer_tf.paragraphs[0]
        footer_p.text = report.title
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = footer_color
        footer_p.alignment = PP_ALIGN.RIGHT

    from tools.report_generate import _safe_filename

    filename = f"{_safe_filename(report.title)}.pptx"
    output_path = output_dir / filename
    prs.save(str(output_path))
    return f"PPT报告已生成(可编辑): {output_path}"


# ---------------------------------------------------------------------------
# 公开入口
# ---------------------------------------------------------------------------


async def render_ppt(report: ReportContent, output_dir: Path, data_context: DataContext | None = None) -> str:
    """PPT 渲染：根据 editable 字段选择可编辑原生模式或图片模式。"""
    if report.editable:
        return _build_native_pptx(report, output_dir, data_context)
    # 图片模式：HTML→Chrome截图→pptx
    pages = _paginate(report)
    html_pages = _render_html_pages(pages, report.template_name)
    images = await _screenshot_pages(html_pages)
    return _build_pptx(images, report.title, output_dir)


async def render_pdf(report: ReportContent, output_dir: Path) -> str:
    """HTML→Chrome直接导出PDF。"""
    pages = _paginate(report)
    html_pages = _render_html_pages(pages, report.template_name)
    return await _export_pdf(html_pages, report.title, output_dir)